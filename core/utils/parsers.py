import os
from typing import Tuple, Dict
import pdfplumber
import pandas as pd
import pptx
import docx
import re

import re
from typing import Tuple, Dict, List, Any


def parse_markdown(path: str) -> Tuple[str, Dict]:
    with open(path, encoding='utf-8') as f:
        content = f.read()

    lines = content.split('\n')
    headings = []
    stack = []  # 用于跟踪标题层级
    current = {"level": 0, "title": "Document Root", "content": "", "elements": [], "children": []}

    # 状态变量
    in_code_block = False
    current_code_block = {"type": "code", "language": "", "content": ""}
    in_table = False
    current_table = {"type": "table", "headers": [], "rows": []}

    for line in lines:
        # 检测标题行
        match = re.match(r'(#+) (.+)', line)
        if match:
            # 结束当前代码块或表格状态
            if in_code_block:
                current["elements"].append(current_code_block)
                in_code_block = False
                current_code_block = {"type": "code", "language": "", "content": ""}
            if in_table:
                current["elements"].append(current_table)
                in_table = False
                current_table = {"type": "table", "headers": [], "rows": []}

            level = len(match.group(1))
            title = match.group(2)

            # 创建新标题节点
            new_heading = {
                "level": level,
                "title": title,
                "content": "",
                "elements": [],
                "children": []
            }

            # 确定新标题在层级结构中的位置
            while stack and stack[-1]["level"] >= level:
                stack.pop()

            if stack:
                # 添加到父标题的子节点
                stack[-1]["children"].append(new_heading)
            else:
                # 顶级标题
                headings.append(new_heading)

            stack.append(new_heading)
            current = new_heading
            continue

        # 处理代码块
        code_match = re.match(r'^```(\w*)$', line)
        if code_match:
            if in_code_block:
                # 结束代码块
                current["elements"].append(current_code_block)
                in_code_block = False
                current_code_block = {"type": "code", "language": "", "content": ""}
            else:
                # 开始代码块
                in_code_block = True
                current_code_block["language"] = code_match.group(1)
            continue

        if in_code_block:
            # 在代码块中，收集内容
            current_code_block["content"] += line + "\n"
            continue

        # 处理表格
        if re.match(r'^\|(.+)\|$', line) and '---' not in line:
            if not in_table:
                # 开始表格
                in_table = True
                # 提取表头
                headers = [h.strip() for h in line.split('|')[1:-1]]
                current_table["headers"] = headers
            else:
                # 表格行
                row = [cell.strip() for cell in line.split('|')[1:-1]]
                current_table["rows"].append(row)
            continue

        # 处理表格分隔线（表示表头结束）
        if re.match(r'^\|(\s*:?---+:?\s*\|)+$', line):
            continue  # 跳过表格分隔线

        # 结束表格状态
        if in_table and not re.match(r'^\|', line):
            current["elements"].append(current_table)
            in_table = False
            current_table = {"type": "table", "headers": [], "rows": []}

        # 处理图片
        img_match = re.match(r'!\[(.*?)\]\((.*?)\)', line)
        if img_match:
            alt_text = img_match.group(1)
            img_url = img_match.group(2)
            current["elements"].append({
                "type": "image",
                "alt": alt_text,
                "url": img_url
            })
            continue

        # 处理链接
        link_match = re.match(r'\[(.*?)\]\((.*?)\)', line)
        if link_match:
            link_text = link_match.group(1)
            link_url = link_match.group(2)
            current["elements"].append({
                "type": "link",
                "text": link_text,
                "url": link_url
            })
            continue

        # 处理列表
        list_match = re.match(r'^(\s*)([-*]|\d+\.)\s+(.*)', line)
        if list_match:
            indent = len(list_match.group(1))
            list_type = "unordered" if list_match.group(2) in ['-', '*'] else "ordered"
            item_text = list_match.group(3)

            # 添加到当前元素
            if current["elements"] and current["elements"][-1]["type"] == "list":
                # 继续当前列表
                current_list = current["elements"][-1]
            else:
                # 开始新列表
                current_list = {
                    "type": "list",
                    "list_type": list_type,
                    "items": []
                }
                current["elements"].append(current_list)

            current_list["items"].append({
                "indent": indent,
                "text": item_text
            })
            continue

        # 普通文本行
        if stack:
            stack[-1]["content"] += line + "\n"
        else:
            current["content"] += line + "\n"

    # 处理结束时可能仍在代码块或表格中的情况
    if in_code_block:
        current["elements"].append(current_code_block)
    if in_table:
        current["elements"].append(current_table)

    # 添加文档结构信息
    doc_structure = {
        "type": "markdown",
        "content": content,
        "headings": headings
    }

    return content, doc_structure


import fitz  # PyMuPDF
import pdfplumber
import re
from typing import Tuple, Dict, List, Any


def parse_pdf(path: str) -> Tuple[str, Dict]:
    # 使用PyMuPDF提取结构化元素
    doc = fitz.open(path)
    structured_data = {
        "type": "pdf",
        "pages": len(doc),
        "elements": []
    }

    # 提取文档元数据
    metadata = doc.metadata
    structured_data["metadata"] = {
        "title": metadata.get("title", ""),
        "author": metadata.get("author", ""),
        "subject": metadata.get("subject", ""),
        "keywords": metadata.get("keywords", ""),
        "creation_date": metadata.get("creationDate", ""),
        "modification_date": metadata.get("modDate", "")
    }

    full_text = ""

    with pdfplumber.open(path) as pdf:
        for page_num in range(len(doc)):
            page_ = pdf.pages[page_num]
            page = doc.load_page(page_num)
            page_text = f"\n\n--- Page {page_num + 1} ---\n\n"

            # 获取页面中的所有内容块（按顺序）
            blocks = page.get_text("blocks", sort=True)

            # 提取表格和公式（用于添加到full_text）
            tables = page_.extract_tables()
            formulas = re.findall(r'\$(.*?)\$', page.get_text("text"))

            for block in blocks:
                # 块类型: 0=文本, 1=图像, 2=路径, 3=注释, 4=表单域
                if block[6] == 0:  # 文本块
                    block_text = block[4].strip()
                    if block_text:
                        page_text += block_text + "\n"
                elif block[6] == 1:  # 图像块
                    page_text += f"[Image: {block[7]} at position ({block[0]}, {block[1]})]\n"

            # 添加表格内容到full_text
            if tables:
                page_text += "\n[表格内容]\n"
                for i, table in enumerate(tables):
                    page_text += f"表格 {i + 1}:\n"
                    for row in table:
                        # 清理表格内容并转换为文本行
                        cleaned_row = [cell.replace("\n", " ") if cell else "" for cell in row]
                        page_text += "| " + " | ".join(cleaned_row) + " |\n"
                    page_text += "\n"

            # 添加公式内容到full_text
            if formulas:
                page_text += "\n[公式内容]\n"
                for i, formula in enumerate(formulas):
                    page_text += f"公式 {i + 1}: ${formula}$\n"

            full_text += page_text

            # 提取当前页的结构化元素
            page_elements = extract_page_elements(doc, page_num, page_)
            structured_data["elements"].extend(page_elements)

    return full_text, structured_data


def extract_page_elements(doc: fitz.Document, page_num: int, page_) -> List[Dict[str, Any]]:
    """提取单页的结构化元素"""
    page = doc.load_page(page_num)
    elements = []

    # 提取文本块（标题和段落）
    blocks = page.get_text("dict")["blocks"]
    for b in blocks:
        if "lines" in b:  # 文本块
            block_text = ""
            for line in b["lines"]:
                for span in line["spans"]:
                    block_text += span["text"] + " "

            # 根据字体大小判断标题级别
            font_size = b["lines"][0]["spans"][0]["size"]
            is_heading = font_size > 12  # 假设标题字体大于12pt

            element = {
                "type": "heading" if is_heading else "paragraph",
                "level": determine_heading_level(font_size),
                "text": block_text.strip(),
                "bbox": b["bbox"],
                "page": page_num + 1
            }
            elements.append(element)

    # 提取图片
    images = page.get_image_info()
    for img in images:
        elements.append({
            "type": "image",
            "bbox": img["bbox"],
            "width": img["width"],
            "height": img["height"],
            "page": page_num + 1
        })

    # 提取表格（使用pdfplumber更准确）
    page_obj = page_
    tables = page_obj.extract_tables()
    for i, table in enumerate(tables):
        # 转换为二维列表
        table_data = []
        for row in table:
            cleaned_row = [cell.replace("\n", " ") if cell else "" for cell in row]
            table_data.append(cleaned_row)

        elements.append({
            "type": "table",
            "data": table_data,
            "page": page_num + 1,
            "table_index": i
        })

    # 尝试识别公式（基于LaTeX模式）
    text = page.get_text("text")
    latex_formulas = re.findall(r'\$(.*?)\$', text)  # 简单匹配$...$格式的公式
    for formula in latex_formulas:
        elements.append({
            "type": "formula",
            "format": "latex",
            "content": formula,
            "page": page_num + 1
        })

    return elements


def determine_heading_level(font_size: float) -> int:
    """根据字体大小确定标题级别"""
    if font_size > 20:
        return 1  # 一级标题
    elif font_size > 16:
        return 2  # 二级标题
    elif font_size > 14:
        return 3  # 三级标题
    elif font_size > 12:
        return 4  # 四级标题
    else:
        return 0  # 普通文本

def parse_excel(path: str) -> Tuple[str, Dict]:
    xls = pd.read_excel(path, sheet_name=None)
    content = []
    structured = {"type": "excel_sheet", "sheets": []}

    for sheet, df in xls.items():
        df = df.fillna("")
        sheet_lines = df.astype(str).values.tolist()
        content.append(f"[Sheet: {sheet}]\n" + df.to_string(index=False))
        structured["sheets"].append({
            "name": sheet,
            "headers": list(df.columns),
            "rows": sheet_lines
        })

    return "\n\n".join(content), structured


def parse_pptx(path: str) -> Tuple[str, Dict]:
    prs = pptx.Presentation(path)
    slides_text = []
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_title = f"Slide {i + 1}"
        slide_content = []
        slide_elements = []

        # 提取幻灯片标题（如果有）
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip()

        # 处理所有形状
        for shape in slide.shapes:
            element = None

            # 文本框处理
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    # 判断是标题还是正文
                    if shape == slide.shapes.title:
                        element = {"type": "title", "text": text}
                    else:
                        # 检查是否有项目符号
                        is_bullet = any([p.level > 0 for p in shape.text_frame.paragraphs])
                        element = {
                            "type": "paragraph",
                            "text": text,
                            "is_bullet": is_bullet
                        }
                    slide_content.append(text)

            # 表格处理
            elif shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)

                element = {
                    "type": "table",
                    "data": table_data
                }
                slide_content.append("表格内容")

            # 图片处理
            elif shape.shape_type == 13:  # 13 表示图片类型
                # 获取图片数据
                try:
                    image = shape.image
                    element = {
                        "type": "image",
                        "width": shape.width,
                        "height": shape.height
                    }
                    slide_content.append("[图片]")
                except Exception as e:
                    print(f"Error processing image in slide {i + 1}: {e}")

            if element:
                slide_elements.append(element)

        # 构建幻灯片文本
        slide_text = f"{slide_title}:\n" + "\n".join(slide_content)
        slides_text.append(slide_text)

        # 构建幻灯片结构化数据
        slides_data.append({
            "title": slide_title,
            "elements": slide_elements
        })

    # 构建最终结果
    full_text = "\n\n".join(slides_text)
    structured_data = {
        "type": "ppt_teaching",
        "slides": slides_data,
        "total_slides": len(prs.slides)
    }

    return full_text, structured_data


def parse_docx(path: str) -> Tuple[str, Dict]:
    doc = docx.Document(path)
    full_text = []
    structured = {
        "type": "word_document",
        "elements": [],
        "sections": []
    }

    current_section = {"title": None, "elements": []}

    # 遍历文档中的所有段落
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 检查标题样式
        if para.style.name.startswith('Heading'):
            level = int(para.style.name.split(' ')[-1]) if ' ' in para.style.name else 1

            # 保存当前章节（如果有内容）
            if current_section["title"] or current_section["elements"]:
                structured["sections"].append(current_section)

            # 开始新章节
            current_section = {
                "title": text,
                "level": level,
                "elements": []
            }
            full_text.append(f"\n\n{'#' * level} {text}\n")

            # 添加到结构化元素
            structured["elements"].append({
                "type": "heading",
                "level": level,
                "text": text
            })
            continue

        # 处理列表
        if para.style.name.startswith('List'):
            list_type = "unordered" if "Bullet" in para.style.name else "ordered"

            # 添加到结构化元素
            structured["elements"].append({
                "type": "list_item",
                "list_type": list_type,
                "text": text
            })
            full_text.append(f"- {text}\n")
            continue

        # 普通段落
        full_text.append(text + "\n")
        structured["elements"].append({
            "type": "paragraph",
            "text": text
        })

        # 添加到当前章节
        if current_section:
            current_section["elements"].append({
                "type": "paragraph",
                "text": text
            })

    # 处理文档中的表格
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            table_data.append(row_data)

        # 添加到结构化元素
        structured["elements"].append({
            "type": "table",
            "data": table_data
        })

        # 添加到当前章节
        if current_section:
            current_section["elements"].append({
                "type": "table",
                "data": table_data
            })

        # 添加到文本表示
        full_text.append("\n[表格开始]\n")
        for row in table_data:
            full_text.append("| " + " | ".join(row) + " |\n")
        full_text.append("[表格结束]\n")

    # 处理文档中的图片
    rels = doc.part.rels
    for rel in rels:
        if "image" in rels[rel].target_ref:
            image_part = rels[rel].target_part
            image_ext = image_part.content_type.split("/")[-1]
            image_name = f"image_{len(structured['elements']) + 1}.{image_ext}"

            # 添加到结构化元素
            structured["elements"].append({
                "type": "image",
                "description": f"图片 {len(structured['elements']) + 1}"
            })

            # 添加到当前章节
            if current_section:
                current_section["elements"].append({
                    "type": "image",
                    "description": f"图片 {len(structured['elements']) + 1}"
                })

            # 添加到文本表示
            full_text.append(f"\n[图片: {image_name}]\n")

    # 添加最后一个章节
    if current_section["title"] or current_section["elements"]:
        structured["sections"].append(current_section)

    return "".join(full_text), structured