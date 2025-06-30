import json
import math
from pathlib import Path
from pprint import pprint
import configparser

import pdfplumber
import yaml
import matplotlib.pyplot as plt
from docx import Document
from pptx import Presentation
from transformers.trainer_callback import TrainerCallback

IGNORE_INDEX = -100

def read_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        pages = [page.extract_text() for page in pdf.pages if page.extract_text()]
    return "\n\n".join(pages)

def read_docx(file_path):
    doc = Document(file_path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paras)

def read_pptx(file_path):
    prs = Presentation(file_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text.strip())
        slides.append(f"Slide {idx+1}:\n" + "\n".join(parts))
    return "\n\n".join(slides)

def read_md(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def read_file(file_path):
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return read_pdf(file_path), "PDF"
    elif suffix == ".docx":
        return read_docx(file_path), "Word"
    elif suffix == ".pptx" or suffix == ".ppt":
        return read_pptx(file_path), "PPT"
    elif suffix == ".md":
        return read_md(file_path), "Markdown"
    else:
        return None, None

def load_config(config_path):
    """Load config from yaml path."""
    with open(config_path, "r") as fp:
        config = yaml.safe_load(fp)
    return config

def load_json(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return json.load(f)

def load_secret(secret_path):
    """Load secret from ini path."""
    config = configparser.ConfigParser()
    config.read(secret_path)
    return config

class LossRecorderCallback(TrainerCallback):
    def __init__(self):
        super().__init__()
        self.training_loss = []
        self.eval_loss = []
        self.steps = []
        self.eval_steps = []

    def on_log(self, args, state, control, logs=None, **kwargs):
        if logs is None:
            return

        # Record training loss
        if "loss" in logs:
            self.training_loss.append(logs["loss"])
            self.steps.append(state.global_step)

        # Record evaluation loss (skipping nan values)
        if "eval_loss" in logs and not math.isnan(logs["eval_loss"]):
            self.eval_loss.append(logs["eval_loss"])
            self.eval_steps.append(state.global_step)


def plot_training_loss(callback, save_path="training_loss_plot.png"):
    """Plot and save the training and evaluation loss curves"""
    plt.figure(figsize=(10, 6))

    # Plot training loss
    if callback.training_loss:
        plt.plot(callback.steps, callback.training_loss, label="Training Loss", color="blue")

    # Plot evaluation loss
    if callback.eval_loss:
        plt.plot(callback.eval_steps, callback.eval_loss, label="Evaluation Loss", color="red", marker="o")

    plt.title("Training and Evaluation Loss")
    plt.xlabel("Steps")
    plt.ylabel("Loss")
    plt.legend()
    plt.grid(True)
    plt.tight_layout()
    plt.savefig(save_path)
    print(f"Loss plot saved to {save_path}")